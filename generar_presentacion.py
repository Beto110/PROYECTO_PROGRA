import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Crear la presentación
prs = Presentation()

# Función para agregar títulos con estilo
def apply_title_style(title_shape, text):
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102) # Azul oscuro

# 1. Diapositiva de Portada
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
apply_title_style(title, "InvenWeb\nSistema Web de Control de Inventario")
subtitle.text = ("Integrantes:\n"
                 "- Alberto Jose Velasquez Paz\n"
                 "- Ricardo Alberto Mendiola Hernández\n"
                 "- Luis Angel Zúniga Menjivar\n"
                 "- Victor Arnoldo Iglesias Sandoval\n"
                 "- Shelsy Rubi Melendez Rodríguez")

# 2. Situación Problemática (Alberto)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
apply_title_style(title, "Situación Problemática")
tf = body.text_frame
tf.text = "El problema del control manual en los negocios"
p = tf.add_paragraph()
p.text = "• Pérdida de información importante."
p.level = 1
p = tf.add_paragraph()
p.text = "• Descuadre entre ventas e inventario real en bodega."
p.level = 1
p = tf.add_paragraph()
p.text = "• Dificultad para tomar decisiones por falta de reportes."
p.level = 1

# 3. Planteamiento y Alcances (Alberto)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
apply_title_style(title, "Planteamiento de la Solución")
tf = body.text_frame
tf.text = "InvenWeb: Plataforma centralizada de gestión"
p = tf.add_paragraph()
p.text = "• Alcance: Registro de entradas, cálculo de salidas (ventas) en tiempo real."
p.level = 1
p = tf.add_paragraph()
p.text = "• Accesible desde el navegador sin instalaciones complejas."
p.level = 1

# 4. Módulos y Arquitectura (Ricardo)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
apply_title_style(title, "Módulos y Arquitectura Técnica")
tf = body.text_frame
tf.text = "Principales Módulos:"
p = tf.add_paragraph()
p.text = "• Categorías, Productos, Usuarios y Carrito de Ventas."
p.level = 1
p = tf.add_paragraph()
p.text = "Desglosamiento Técnico:"
p.level = 0
p = tf.add_paragraph()
p.text = "• Frontend: HTML5, CSS3, JavaScript (Interfaces dinámicas)."
p.level = 1
p = tf.add_paragraph()
p.text = "• Backend: PHP puro procesando lógica de negocio y peticiones."
p.level = 1

# 5. Base de Datos (Luis)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
apply_title_style(title, "Base de Datos y Relaciones")
tf = body.text_frame
tf.text = "Gestor Utilizado: SQLite"
p = tf.add_paragraph()
p.text = "• Ventaja: Despliegue rápido, ligero y embebido."
p.level = 1
p = tf.add_paragraph()
p.text = "Relaciones Principales:"
p.level = 0
p = tf.add_paragraph()
p.text = "• [Aquí insertarán la imagen del Diagrama ER en el PowerPoint final]"
p.level = 1
p = tf.add_paragraph()
p.text = "• Ventas 1:N Detalle_Ventas (Garantiza integridad transaccional)."
p.level = 1

# 6. Validaciones (Victor)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
apply_title_style(title, "Demostración de Validaciones")
tf = body.text_frame
tf.text = "Protegiendo la integridad de los datos"
p = tf.add_paragraph()
p.text = "• Prevención de errores lógicos: Bloqueo de ventas con stock cero."
p.level = 1
p = tf.add_paragraph()
p.text = "• Doble capa de validación:"
p.level = 1
p = tf.add_paragraph()
p.text = "  1. Javascript (Evita enviar formularios vacíos)."
p.level = 2
p = tf.add_paragraph()
p.text = "  2. PHP (Evita inyección y formatos inválidos en backend)."
p.level = 2

# 7. Demo en Vivo (Shelsy)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
apply_title_style(title, "Demostración en Vivo")
tf = body.text_frame
tf.text = "Flujo del Usuario Común"
p = tf.add_paragraph()
p.text = "• 1. Autenticación en el sistema."
p.level = 1
p = tf.add_paragraph()
p.text = "• 2. Agregar productos al inventario."
p.level = 1
p = tf.add_paragraph()
p.text = "• 3. Realizar proceso de venta (Carrito)."
p.level = 1
p = tf.add_paragraph()
p.text = "• 4. Comprobación de reducción de stock."
p.level = 1

# Guardar la presentación
prs.save('Presentacion_InvenWeb.pptx')
print("¡La presentación Presentacion_InvenWeb.pptx ha sido generada exitosamente!")
